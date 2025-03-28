import os
import requests
from bs4 import BeautifulSoup
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 📌 1. 다운로드 폴더 설정
download_dir = "download"
output_text_dir = "extracted_text"
output_image_dir = "extracted_images"
os.makedirs(download_dir, exist_ok=True)
os.makedirs(output_text_dir, exist_ok=True)
os.makedirs(output_image_dir, exist_ok=True)

# 📌 2. 크롤링할 URL 설정
url = "https://www.fustero.es/index_ja.php"
session = requests.Session()  # 세션 사용으로 연결 재사용
response = session.get(url)
response.encoding = 'utf-8'

# 📌 3. 파일 다운로드
def download_file(file_url):
    full_url = requests.compat.urljoin(url, file_url)
    file_name = os.path.join(download_dir, os.path.basename(file_url))

    print(f"📥 다운로드 중: {full_url}")
    file_response = session.get(full_url, stream=True)

    if file_response.status_code == 200:
        with open(file_name, 'wb') as f:
            for chunk in file_response.iter_content(4096):  # 청크 크기 조정
                f.write(chunk)
        print(f"✅ 저장 완료: {file_name}")
    else:
        print(f"❌ 다운로드 실패: {full_url}")

if response.status_code == 200:
    soup = BeautifulSoup(response.text, 'html.parser')
    sections = soup.find_all('section', class_=['w3-container', 'w3-row'])

    file_links = []
    for section in sections:
        for link in section.find_all('a', href=True):
            file_url = link['href']
            if file_url.endswith(".pptx"):  # ✅ PPTX 파일만 다운로드
                file_links.append(file_url)

    # ✅ 멀티스레딩을 활용한 병렬 다운로드
    with ThreadPoolExecutor(max_workers=5) as executor:
        executor.map(download_file, file_links)

else:
    print(f"❌ 페이지 요청 실패: {response.status_code}")

# 📌 4. PPTX에서 텍스트 & 이미지 추출
def extract_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]

    if len(prs.slides) > 1:
        slide = prs.slides[1]  # ✅ 두 번째 슬라이드 추출

        # 📌 4-1. 텍스트 추출
        text_content = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        text_output_path = os.path.join(output_text_dir, f"{base_name}.txt")
        with open(text_output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(text_content))
        print(f"✅ 텍스트 저장 완료: {text_output_path}")

        # 📌 4-2. 이미지 추출
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_bytes = shape.image.blob
                image_ext = shape.image.ext or "png"
                image_output_path = os.path.join(output_image_dir, f"{base_name}_image.{image_ext}")

                with open(image_output_path, "wb") as img_file:
                    img_file.write(image_bytes)
                print(f"✅ 이미지 저장 완료: {image_output_path}")

# ✅ 다운로드한 모든 PPTX 파일에서 텍스트 & 이미지 추출
for file_name in os.listdir(download_dir):
    if file_name.endswith(".pptx") and not file_name.startswith("~$"):
        extract_from_pptx(os.path.join(download_dir, file_name))

print("🎉 모든 작업 완료!")
input("엔터 키를 눌러 종료하세요...")  # 프로그램이 바로 닫히지 않도록 유지
