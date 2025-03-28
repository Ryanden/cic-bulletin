import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 다운로드 폴더 경로
download_dir = "download"

# 저장할 폴더 생성
output_text_dir = "extracted_text"
output_image_dir = "extracted_images"
os.makedirs(output_text_dir, exist_ok=True)
os.makedirs(output_image_dir, exist_ok=True)

# 다운로드 폴더 내의 모든 .pptx 파일 처리
for file_name in os.listdir(download_dir):
    if file_name.endswith(".pptx") and not file_name.startswith("~$"):  # ✅ 임시 파일 무시
        pptx_path = os.path.join(download_dir, file_name)
        prs = Presentation(pptx_path)


        # 슬라이드가 2개 이상 있는지 확인
        if len(prs.slides) > 1:
            slide = prs.slides[1]  # 두 번째 슬라이드 (인덱스는 0부터 시작)

            # 📌 1. 텍스트 추출
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

            base_name = os.path.splitext(file_name)[0]  # 확장자 제거
            text_output_path = os.path.join(output_text_dir, f"{base_name}.txt")
            with open(text_output_path, "w", encoding="utf-8") as f:
                f.write("\n".join(text_content))
            print(f"✅ 텍스트 저장 완료: {text_output_path}")

            # 📌 2. 이미지 추출
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob  # 이미지 데이터

                    # 파일 확장자 가져오기 (JPEG, PNG 등)
                    ext = image.ext or "png"
                    image_output_path = os.path.join(output_image_dir, f"{file_name}_image.{ext}")

                    with open(image_output_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    print(f"✅ 이미지 저장 완료: {image_output_path}")
        else:
            print(f"⚠ {file_name}: 슬라이드가 2개 미만이라 처리하지 않음.")
